import copy
import json
import re
import time
from pathlib import Path

import requests
from lxml import html
from openpyxl import load_workbook
from PIL import ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE


ROOT = Path(__file__).resolve().parent
CORRECTED_XLSX = ROOT / "Обыкновенное чудо - Субтитры с исправлениями от 18 сентября.xlsx"
SOURCE_PPTX = ROOT / "Обыкновенное чудо - Субтитры.pptx"
OUTPUT_XLSX = ROOT / "Обыкновенное чудо - Субтитры - логичная разбивка.xlsx"
OUTPUT_PPTX = ROOT / "Обыкновенное чудо - Субтитры - логичная разбивка.pptx"
CACHE_PATH = ROOT / ".translation-cache.json"
KEEP_ROWS = 300
SUBTITLE_FONT = ImageFont.truetype(r"C:\Windows\Fonts\calibril.ttf", 75)
# The box is 1184 px wide at 96 DPI. The safety margin absorbs PowerPoint's
# glyph and word-spacing differences from Pillow.
MAX_LINE_WIDTH = 1100

EDITORIAL_OVERRIDES = {
    302: (None, "How can magic appear out of nowhere? I feel so happy, though I don't know why."),
    304: (None, "It feels like a holiday. Perhaps New Year's Day? How else can I explain..."),
    305: (None, "why listening to you and looking at you feels so wonderful? Why do I feel so warm?"),
    306: (None, "Such things do not happen... Yet it has just happened!"),
    309: (None, "Why am I so confused? What has come over me? Such strange weakness..."),
    310: (None, "unease and vulnerability. This is so unlike me—as if I were someone else."),
    311: (None, "You would never hurt me, would you?"),
    317: (None, "Those dreary, familiar faces! Shall we run away from them?"),
    354: (None, "CHORUS: I am merely a nobody, less than zero. No one could be more insignificant."),
    364: (None, "Tyrant! Despot! Executioner! Scoundrel! No creature is more repulsive!"),
    374: (None, "Expect no kindness from that parasite! What use is my head..."),
    429: (None, "A woman is wonderful—a lady, signora, femme. A woman's greatest joy..."),
    430: (None, "is having a man beside her, especially here, where nature is wild and innocent."),
    431: (None, "CHORUS: A butterfly flaps—flip-flap-flap; a sparrow follows—hop-hop-hop."),
    432: ("Он её, голубушку – шмяк-шмяк-шмяк, Ам-ням-ням, да и шмыг-шмыг-шмыг!", "He snaps her up—smack-smack-smack, yum-yum-yum—and off he darts!"),
    433: (None, "A butterfly flaps—flip-flap-flap; a sparrow follows—hop-hop-hop."),
    434: ("Он её, голубушку – Шмяк-шмяк-шмяк, Ам-ням-ням, да и Шмыг-шмыг-шмыг!", "He snaps her up—smack-smack-smack, yum-yum-yum—and off he darts!"),
    435: (None, "You are my angel, my ideal, my star, my sweet, my darling. Pearl teeth, coral lips..."),
    436: (None, "a lovely bosom and smile. I've never met your like—let us correct that oversight."),
    437: ("А бабочка крылышками бяк-бяк-бяк,", "A butterfly flaps—flip-flap-flap;"),
    438: (None, "a sparrow follows—hop-hop-hop. He snaps her up—smack-smack-smack,"),
    439: (None, "yum-yum-yum—and off he darts!"),
    494: (None, "I will never, ever let anyone sleep anywhere!"),
    582: (None, "I lost everything just as I found it. I cannot understand this blow from nowhere."),
    612: (None, "Innkeeper, I am... what do you call it? A king."),
    817: (None, "My dear sir, I really ought to go home, but I cannot disappoint the flowers."),
    819: (None, "And I like you even more."),
    821: (None, "Tra-la-la. How delightful and amusing to know that you are fond of me."),
    822: (None, "And I like you even more."),
    824: (None, "Tra-la-la."),
    931: (None, "What sort of burden is that?! It is... it is..."),
    940: (None, "What are you saying? Love?! What love? There is no love!"),
    1114: ("She will kiss him, touch him with her lips, Yet nothing, nothing in the world will change! Farewell!..", "She will kiss him, touch him with her lips, yet nothing in the world will change! Farewell!"),
    1180: (None, "And the lightest down can knock you over. What is happening?"),
    1202: (None, "No, he hasn't. The Hunter is here. The Bear is nowhere to be found."),
    1204: (None, "It isn't the Administrator keeping him out."),
    1209: (None, "Everything will be fine. It will end sadly. Call his friends to say goodbye."),
    1214: (None, "Hello. This is my wife, Amanda. We've had a baby girl."),
    1220: (None, "— What greatness?\n— It makes those left alive reflect."),
    1222: (None, "It is shameful to kill heroes to stir the cold-hearted and move the indifferent."),
    1225: (None, "He's playing cards. The old gadabout."),
    1232: (None, "No need to answer. It is obvious that all is well. Everything is splendid."),
    1238: (None, "Hush-a-bye, come closer. Prop up my roof—it has slipped. I know nothing. I see nothing."),
    1239: (None, "I hear nothing. Well, how about that!"),
    1240: (None, "Boyish! Childish! You'll sleep through it all. At least try to help your daughter!"),
    1241: (None, "Stop the Prince! Dismiss the rogue! He'll be the death of your country and daughter."),
    1242: (None, "All day his agents haul bales and boxes back and forth, heedless of everything."),
    1245: (None, "Hush-a-bye, come closer. Prop up my roof—it has slipped. I know nothing. I see nothing."),
    1246: (None, "I hear nothing. Well, how about that!"),
    1249: (None, "Because I am degenerating, you fool!"),
    1250: (None, "Read books instead of demanding what the King cannot do! What?!"),
    1251: (None, "The Princess will die? So let her!"),
    1253: (None, "I recently got some poison and tried it on the cook. He dropped dead without even noticing."),
    1254: (None, "Then allow me to tell you: You, Your Excellency..."),
    1255: (None, "are, in my opinion, a fool!"),
    1259: (None, "To call me, the greatest of kings, by a general's title! Is this a revolt?"),
    1265: (None, "How do you like that? And I'll go further: rumors of your holiness are exaggerated!"),
    1266: (None, "You do not deserve to be an Honorary Saint!"),
    1270: (None, "An honorary pope? You're no pope, no pope, no pope!"),
    1271: (None, "I am a father... This is too much. Executioner!"),
    1276: (None, "Who dares offend our splendid, regular-guy little King, as I call him?"),
    1281: (None, "What nonsense! Sheer delirium, as I call it."),
    1282: (None, "The doctor the King and I share examined her yesterday and found no illness caused by love."),
    1284: (None, "First. Second, love causes amusing, entirely curable ailments."),
    1287: (None, "The doctor staked his head on the Princess recovering any moment."),
    1290: (None, "Disaster! Disaster! The doctor has fled!"),
    1293: (None, "I went for some calming drops and found the rooms unlocked, drawers open, cupboards empty—and a note."),
    1294: (None, "Here it is!"),
    1296: (None, "They took away my gendarmes. Now they frighten me. You are pigs, not loyal subjects."),
    1298: (None, "Yes... Our little King has grown old."),
    1299: (None, "Anyone would grow old with you around, damn it..."),
    1303: (None, "Then you'll blame me. A doctor is only human, with human weaknesses. He wants to live."),
    1305: (None, "Why are you standing there? Bring back the doctor! Blame him for everything! Move!"),
    1308: (None, "How wonderful—you are here too, my friend Wizard. And you. What an extraordinary day!"),
    1311: (None, "Life is smiling on me one last time. Have they told you I am to die today?"),
    1314: (None, "My friends, be even kinder to me than ever."),
    1318: (None, "Anything to keep me from thinking of what will soon happen to me. Orinthia..."),
    1320: (None, "Not as happy as we imagined, but happy."),
    1325: (None, "I don't know, Princess. I think I'm not bad."),
    1326: (None, "I just love my husband and child so very much..."),
    1329: (None, "But I did nothing but pine. Life... The same people, the same words."),
    1331: (None, "One dull note keeps creaking, aching: He still isn't here, still isn't here."),
    1332: (None, "No use waiting. Seconds pass; my life drains away. Tonight I know for sure..."),
    1333: (None, "It will run out. A black-clad crone with a scythe will cloak me, and I will..."),
    1334: (None, "follow her meekly."),
    1339: (None, "No. It's him. It's him!"),
    1341: (None, "— Have you come to see me?\n— Yes. Hello. Why are you crying?"),
    1342: (None, "They are tears of joy."),
    1344: (None, "Only you. Here is my secret: I love you! It is true. I love you enough to forgive anything!"),
    1345: (None, "You want to turn into a bear? So be it! Just don't leave me!"),
    1346: (None, "I can't go on wasting away here alone! Why were you gone so long?"),
    1348: (None, "If you didn't come, you couldn't. See how docile I've become? Just don't leave!"),
    1350: (None, "I was coming. I kept coming all this time. Don't be angry. I'm here."),
    1360: (None, "Yes."),
    1366: (None, "Some performed better, some worse, but by now I've grown fond of them."),
    1367: (None, "I can't cross them out! They're people, not words. Sleep, my darling. Let it be."),
    1370: (None, "Perhaps you won't die, but become ivy and twine around this fool."),
    1372: (None, "You're angry... But look what I've thought of. Sleep."),
    1373: (None, "You may wake to find tomorrow here and every sorrow left in yesterday. Sleep, my dear."),
    1374: ("— Вы горюете, друзья?\n— Да.", "— Are you grieving, my friends?\n— Yes."),
    1378: (None, "A young man may kiss a Princess and not become a Bear—or no one notices if he does!"),
    1380: (None, "Still, forgive me for building castles in the air."),
    1382: (None, "Rain is one thing, but there are miracles, wondrous changes, and comforting dreams."),
    1386: (None, "You may wake to find tomorrow here and every sorrow left in yesterday. Sleep."),
    1388: ("Я дал слово. Тихо!", "I gave my word. Be quiet!"),
    1390: (None, "Frightened? Don't be. Let's go back to your room."),
    1391: (None, "Look, everyone is asleep: the sentries on the towers, my father on his throne."),
    1395: (None, "Suddenly we're alone in all the world. Wait..."),
    1398: (None, "Oh God, I'm so happy I decided as I did."),
    1399: (None, "Fool that I was, I never knew how wonderful this could be."),
    1403: ("Пусть будет так, как ты хочешь. Пусть будет так, как ты хочешь.", "Let it be as you wish. Let it be as you wish."),
    1407: (None, "Love transformed him so completely that he can never be a bear again."),
    1413: (None, "Let us speak quietly, in low voices."),
    1414: (None, "Let's part with light hearts. In a week or two, we'll recover. What's past is past."),
    1415: (None, "Of course it's dreadful, absurd, senseless. If only we could bring back the beginning!"),
    1416: (None, "We cannot bring back the beginning. Don't even try. Forget it."),
    1417: (None, "There were tears and laughter; now they're gone. But the echo remains—and that is good."),
    1418: (None, "It came and went."),
}

ROLES = {
    "ВОЛШЕБНИК", "ЖЕНА ВОЛШЕБНИКА", "МЕДВЕДЬ", "ПРИНЦЕССА", "КОРОЛЬ",
    "МИНИСТР-АДМИНИСТРАТОР", "ПЕРВЫЙ МИНИСТР", "ЭМИЛИЯ", "АМАНДА",
    "ОРИНТИЯ", "ЭМИЛЬ", "ОХОТНИК", "УЧЕНИК", "ПАЛАЧ", "ХОР", "ТРАКТИРЩИК",
}
ALIASES = {
    "ЖЕНА": "ЖЕНА ВОЛШЕБНИКА",
    "ВОЛШЕБНИЦА": "ЖЕНА ВОЛШЕБНИКА",
    "УЧЕНИК ОХОТНИКА": "УЧЕНИК",
    "МИНИСТР": "МИНИСТР-АДМИНИСТРАТОР",
    "АДМИНИСТРАТОР": "МИНИСТР-АДМИНИСТРАТОР",
}
SERVICE_ROWS = {"Первый звонок", "Второй звонок", "Третий звонок", "Антракт — 20 минут", "Конец"}


def normalize(text):
    return re.sub(r"[^0-9A-Za-zА-Яа-я]+", "", text.replace("ё", "е").replace("Ё", "Е")).lower()


def role_name(text):
    name = re.sub(r"\s*-\s*", "-", re.sub(r"\s+", " ", text.strip().upper()))
    return ALIASES.get(name, name) if name in ROLES or name in ALIASES else None


def detect_role(text):
    match = re.match(r"^([А-ЯЁ][А-ЯЁ \-]{1,40}?)[.:]\s*(.*)$", text, re.I)
    if not match:
        return None
    role = role_name(match.group(1))
    return (role, match.group(2)) if role else None


def clean_spoken(text):
    text = re.sub(r"🎶.*$", "", text)
    text = re.sub(r"\(+[^)]*\)+", " ", text)
    text = re.sub(r"\s+([,.;:!?])", r"\1", text)
    return re.sub(r"\s+", " ", text).strip(" -\t:)")


def performer_role(text):
    found = []
    upper = re.sub(r"\s*-\s*", "-", text.upper())
    for candidate in sorted(ROLES | set(ALIASES), key=len, reverse=True):
        if re.search(rf"\b{re.escape(candidate)}\b", upper):
            mapped = ALIASES.get(candidate, candidate)
            if mapped not in found:
                found.append(mapped)
    if not found:
        return None
    return found[0] if len(found) == 1 else "MIXED"


def script_units():
    document = html.parse(str(ROOT / "script.html"))
    css = "\n".join(document.xpath("//style/text()"))
    italic_classes = {
        name for name, body in re.findall(r"\.(c\d+)\s*\{([^}]*)\}", css)
        if re.search(r"font-style\s*:\s*italic", body, re.I)
    }
    song_classes = {
        name for name, body in re.findall(r"\.(c\d+)\s*\{([^}]*)\}", css)
        if re.search(r"color\s*:\s*#156082", body, re.I)
    } - italic_classes

    units = []
    current_role = None
    pending_role = None
    last_was_pause = False

    def add_pause():
        nonlocal last_was_pause
        if units and not last_was_pause:
            units.append({"pause": True, "ru": "***", "role": "MIXED", "song": False})
            last_was_pause = True

    def add_text(text, role, song):
        nonlocal last_was_pause
        cleaned = clean_spoken(text)
        if cleaned:
            units.append({"pause": False, "ru": cleaned, "role": role or "MIXED", "song": song})
            last_was_pause = False

    for element in document.xpath("//p|//h1|//h2|//h3|//h4|//li"):
        classes = set(element.get("class", "").split())
        for span in element.xpath(".//span"):
            if set(span.get("class", "").split()) & italic_classes:
                span.drop_tree()
        text = re.sub(r"\s+", " ", element.text_content().replace("\xa0", " ")).strip()
        if not text:
            continue
        if element.tag in {"h1", "h2", "h3", "h4"} or re.match(r"^СЦЕНА\s+\d+\.\d+", text):
            add_pause()
            current_role = None
            pending_role = None
            continue
        if text.startswith("🎶"):
            continue
        if "c16" in classes:
            pending_role = performer_role(text)
            current_role = None
            continue

        found = detect_role(text)
        if found:
            current_role, remainder = found
            pending_role = None
            add_text(remainder, current_role, bool(classes & song_classes))
            continue

        song = bool(classes & song_classes) or any(
            set(span.get("class", "").split()) & song_classes for span in element.xpath(".//span")
        )
        if song and not current_role and pending_role:
            current_role = pending_role
            pending_role = None
        if current_role:
            inline = re.search(r"\bВДВОЕМ\s*[.:]\s*", text, re.I)
            if inline:
                add_text(text[:inline.start()], current_role, song)
                current_role = "MIXED"
                add_text(text[inline.end():], current_role, song)
            else:
                add_text(text, current_role, song)
    return units


def split_text(text):
    pieces = re.split(r"(?<=[.!?…])\s+|(?<=,)\s+(?=[А-ЯЁ])", text)
    return [piece.strip() for piece in pieces if piece.strip()]


def compact_units(units):
    chunks = []
    for unit in units:
        if unit["pause"]:
            if chunks and not chunks[-1]["pause"]:
                chunks.append(unit.copy())
            continue
        pieces = split_text(unit["ru"])
        current = ""
        for piece in pieces:
            candidate = f"{current} {piece}".strip()
            if current and (len(candidate) > 86 or len(candidate.split()) > 16):
                chunks.append({**unit, "ru": current})
                current = piece
            else:
                current = candidate
        if current:
            chunks.append({**unit, "ru": current})

    merged = []
    for chunk in chunks:
        if chunk["pause"]:
            if merged and not merged[-1]["pause"]:
                merged.append(chunk)
            continue
        if merged and not merged[-1]["pause"]:
            previous = merged[-1]
            candidate = f"{previous['ru']} {chunk['ru']}"
            limit = 94 if previous["song"] and chunk["song"] else 78
            if previous["role"] == chunk["role"] and previous["song"] == chunk["song"] and len(candidate) <= limit and len(candidate.split()) <= 17:
                previous["ru"] = candidate
                continue
        merged.append(chunk)

    result = []
    index = 0
    while index < len(merged):
        row = merged[index]
        if row["pause"]:
            result.append(row)
            index += 1
            continue
        if index + 1 < len(merged):
            following = merged[index + 1]
            if (
                not following["pause"]
                and row["role"] not in {"MIXED", following["role"]}
                and following["role"] != "MIXED"
                and row["ru"].endswith("?")
                and len(row["ru"]) + len(following["ru"]) <= 62
            ):
                result.append({
                    "pause": False,
                    "ru": f"— {row['ru']}\n— {following['ru']}",
                    "role": "MIXED",
                    "song": False,
                })
                index += 2
                continue
        result.append(row)
        index += 1
    return result


def corrected_prefix():
    sheet = load_workbook(CORRECTED_XLSX, read_only=True, data_only=True).active
    return [
        {"ru": str(ru), "en": str(en), "pause": ru == "***"}
        for _, ru, en in sheet.iter_rows(min_row=1, max_row=KEEP_ROWS, min_col=1, max_col=3, values_only=True)
    ]


def trim_to_boundary(units, boundary):
    target = normalize(boundary)
    for index, unit in enumerate(units):
        if unit["pause"]:
            continue
        source = normalize(unit["ru"])
        position = source.find(target)
        if position < 0:
            continue
        normalized_count = 0
        raw_end = 0
        for raw_end, char in enumerate(unit["ru"], 1):
            if normalize(char):
                normalized_count += 1
            if normalized_count >= position + len(target):
                break
        remainder = unit["ru"][raw_end:].strip(" ,")
        tail = []
        if remainder:
            tail.append({**unit, "ru": remainder})
        return tail + units[index + 1:]
    raise ValueError(f"Could not find corrected boundary in script: {boundary}")


def translate_text(text):
    for attempt in range(6):
        try:
            response = requests.get(
                "https://translate.googleapis.com/translate_a/single",
                params={"client": "gtx", "sl": "ru", "tl": "en", "dt": "t", "q": text},
                timeout=30,
            )
            response.raise_for_status()
            return "".join(part[0] for part in response.json()[0]).strip()
        except requests.RequestException:
            if attempt == 5:
                raise
            time.sleep(0.5 * (attempt + 1))


def translate_rows(rows):
    cache = json.loads(CACHE_PATH.read_text(encoding="utf-8")) if CACHE_PATH.exists() else {}
    pending = [index for index, row in enumerate(rows) if not row["pause"]]
    for start in range(0, len(pending), 8):
        if start % 80 == 0:
            print(f"Translation progress: {start}/{len(pending)}")
        indexes = pending[start:start + 8]
        missing = [index for index in indexes if rows[index]["ru"] not in cache]
        if missing:
            marked = "\n".join(f"[{number + 1}] {rows[index]['ru']}" for number, index in enumerate(indexes))
            translated = translate_text(marked)
            matches = list(re.finditer(r"\[(\d+)\]\s*", translated))
            if len(matches) != len(indexes):
                for index in missing:
                    cache[rows[index]["ru"]] = translate_text(rows[index]["ru"])
                    time.sleep(0.05)
            else:
                for number, match in enumerate(matches):
                    end = matches[number + 1].start() if number + 1 < len(matches) else len(translated)
                    cache[rows[indexes[number]]["ru"]] = translated[match.end():end].strip()
            CACHE_PATH.write_text(json.dumps(cache, ensure_ascii=False, indent=2), encoding="utf-8")
            time.sleep(0.08)
        for index in indexes:
            rows[index]["en"] = cache[rows[index]["ru"]]
    return rows


def clone_rgb(color):
    return RGBColor(color[0], color[1], color[2])


def wrapped_lines(text):
    lines = []
    for explicit_line in text.splitlines() or [text]:
        current = ""
        for word in explicit_line.split():
            candidate = f"{current} {word}".strip()
            if current and SUBTITLE_FONT.getlength(candidate) > MAX_LINE_WIDTH:
                lines.append(current)
                current = word
            else:
                current = candidate
        if current:
            lines.append(current)
    return lines


def split_candidates(text):
    normalized = re.sub(r"\s*\n\s*", " ", text).strip()
    normalized = re.sub(r"\bas\s+if\b", "as_if", normalized, flags=re.IGNORECASE)
    normalized = re.sub(r"\bas\s+though\b", "as_though", normalized, flags=re.IGNORECASE)
    clause_start = (
        r"and(?!\s+forth\b)|but|or|that(?!\s+(?:very|same)\b)|when|while|because|"
        r"although|though|as_if|as_though|if|unless|until|"
        r"who|whom|whose|which|where|before|after|"
        r"и|а|но|или|что|чтобы|когда|пока|потому|хотя|если|котор\w*|где|прежде|после|так"
    )
    relative_with_preposition = r"(?:о|об|в|на|с|к|у|для|из|от|по)\s+котор\w*"
    parts = re.split(
        rf"(?<=[.!?…])\s+|(?<=[,;:])\s+|\s+(?=(?:{relative_with_preposition}|{clause_start})\b)",
        normalized,
        flags=re.IGNORECASE,
    )
    return [part.replace("as_if", "as if").replace("as_though", "as though").strip() for part in parts if part.strip()]


def pack_parts(parts):
    chunks = []
    current = ""
    for part in parts:
        candidate = f"{current} {part}".strip()
        if current and len(wrapped_lines(candidate)) > 2:
            chunks.append(current)
            current = part
        else:
            current = candidate
    if current:
        chunks.append(current)
    return chunks


def split_words_balanced(text):
    words = text.split()
    required = max(2, (len(wrapped_lines(text)) + 1) // 2)
    dangling = {
        "a", "an", "the", "to", "of", "for", "from", "with", "without",
        "before", "after", "in", "on", "at", "by", "into", "through", "about",
        "as", "and", "but", "or", "that", "when", "while", "because", "if",
        "is", "are", "was", "were", "be", "been", "being", "has", "have", "had",
        "do", "does", "did", "will", "would", "can", "could", "shall", "should",
    }
    clause_starts = {
        "and", "but", "or", "that", "when", "while", "because", "although",
        "though", "if", "unless", "until", "who", "whom", "whose", "which", "where",
    }
    chunks = []
    start = 0
    for chunk_number in range(1, required):
        remaining_chunks = required - chunk_number
        target_width = SUBTITLE_FONT.getlength(" ".join(words[start:])) / (remaining_chunks + 1)
        candidates = range(start + 1, len(words) - remaining_chunks + 1)

        def boundary_score(end):
            chunk = " ".join(words[start:end])
            last = re.sub(r"[^a-z']", "", words[end - 1].lower())
            next_word = re.sub(r"[^a-z']", "", words[end].lower()) if end < len(words) else ""
            semantic_boundary = bool(re.search(r"[.!?…][\"']?$", words[end - 1])) or next_word in clause_starts
            return (
                max(0, len(wrapped_lines(chunk)) - 2),
                last in dangling or last.endswith("'s"),
                not semantic_boundary,
                abs(SUBTITLE_FONT.getlength(chunk) - target_width),
            )

        end = min(candidates, key=boundary_score)
        chunks.append(" ".join(words[start:end]))
        start = end
    chunks.append(" ".join(words[start:]))
    result = []
    for chunk in chunks:
        if len(wrapped_lines(chunk)) > 2:
            result.extend(split_words_balanced(chunk))
        else:
            result.append(chunk)
    return result


def split_to_fit(text):
    sentences = [part.strip() for part in re.split(r"(?<=[.!?…])\s+", re.sub(r"\s*\n\s*", " ", text)) if part.strip()]
    expanded = []
    for sentence in sentences:
        if len(wrapped_lines(sentence)) <= 2:
            expanded.append(sentence)
            continue
        clauses = split_candidates(sentence)
        if len(clauses) == 1:
            expanded.extend(split_words_balanced(sentence))
        else:
            for clause in clauses:
                if len(wrapped_lines(clause)) > 2:
                    expanded.extend(split_words_balanced(clause))
                else:
                    expanded.append(clause)
    chunks = pack_parts(expanded)

    result = []
    for chunk in chunks:
        if len(wrapped_lines(chunk)) <= 2:
            result.append(chunk)
            continue
        result.extend(split_words_balanced(chunk))
    return result


def split_into_count(text, english_chunks):
    count = len(english_chunks)
    if count == 1:
        return [text]
    parts = split_candidates(text)
    if len(parts) < count:
        parts = text.replace("\n", " ").split()
    lengths = [len(normalize(part)) for part in parts]
    total = sum(lengths)
    english_lengths = [max(1, len(normalize(chunk))) for chunk in english_chunks]
    english_total = sum(english_lengths)
    chunks = []
    start = 0
    for chunk_number in range(1, count):
        target = total * sum(english_lengths[:chunk_number]) / english_total
        minimum_end = start + 1
        maximum_end = len(parts) - (count - chunk_number)
        end = min(
            range(minimum_end, maximum_end + 1),
            key=lambda candidate: abs(sum(lengths[:candidate]) - target),
        )
        chunks.append(" ".join(parts[start:end]).strip())
        start = end
    chunks.append(" ".join(parts[start:]).strip())
    return chunks


def split_overflowing_rows(rows):
    result = []
    split_count = 0
    for row in rows:
        if row["pause"]:
            result.append(row)
            continue
        row["ru"] = re.sub(r"🎵\s*\d+", "", row["ru"]).strip()
        row["en"] = re.sub(r"🎵\s*\d+", "", row["en"]).strip()
        if len(wrapped_lines(row["en"])) <= 2:
            result.append(row)
            continue
        english_chunks = split_to_fit(row["en"])
        russian_chunks = split_into_count(row["ru"], english_chunks)
        for russian, english in zip(russian_chunks, english_chunks):
            result.append({"ru": russian, "en": english, "pause": False})
        split_count += 1
    print(f"Split {split_count} overflowing rows into {len(result) - len(rows) + split_count} rows")
    return result


def reflow_dangling_boundaries(rows):
    russian_starts = {
        "and": r"(?:и|а)\b",
        "but": r"(?:но|а)\b",
        "when": r"когда\b",
        "that": r"(?:что|чтобы)\b",
        "because": r"потому\b",
        "if": r"если\b",
    }
    moved = 0
    for current, following in zip(rows, rows[1:]):
        if current["pause"] or following["pause"]:
            continue
        match = re.search(r"\b(and|but|when|that|because|if)[,;:]?$", current["en"], re.IGNORECASE)
        if not match:
            continue
        conjunction = match.group(1).lower()
        if not re.match(russian_starts[conjunction], following["ru"], re.IGNORECASE):
            continue
        candidate = f"{match.group(1)} {following['en']}"
        if len(wrapped_lines(candidate)) > 2:
            continue
        current["en"] = current["en"][:match.start()].rstrip(" ,;:")
        following["en"] = candidate
        moved += 1
    print(f"Moved {moved} dangling conjunctions to their clauses")
    return rows


def remove_all_slides(presentation):
    slide_ids = presentation.slides._sldIdLst
    for slide_id in list(slide_ids):
        presentation.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def write_outputs(rows):
    workbook = load_workbook(CORRECTED_XLSX)
    sheet = workbook.active
    russian_style = copy.copy(sheet["B1"]._style)
    english_style = copy.copy(sheet["C1"]._style)
    sheet.delete_rows(1, sheet.max_row)
    for number, row in enumerate(rows, 1):
        english = row.get("en", "***" if row["pause"] else "")
        sheet.cell(number, 2, row["ru"])._style = copy.copy(russian_style)
        sheet.cell(number, 3, english)._style = copy.copy(english_style)
        if "\n" in row["ru"] or "\n" in english:
            sheet.row_dimensions[number].height = 31
    workbook.save(OUTPUT_XLSX)

    presentation = Presentation(SOURCE_PPTX)
    sample = next(shape for shape in presentation.slides[0].shapes if shape.has_text_frame)
    sample_run = next(run for paragraph in sample.text_frame.paragraphs for run in paragraph.runs)
    background_rgb = presentation.slides[0].background.fill.fore_color.rgb
    remove_all_slides(presentation)
    layout = min(presentation.slide_layouts, key=lambda item: len(item.placeholders))
    for row in rows:
        slide = presentation.slides.add_slide(layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = clone_rgb(background_rgb)
        if row["pause"]:
            continue
        box = slide.shapes.add_textbox(sample.left, sample.top, presentation.slide_width - 2 * sample.left, 1518285)
        frame = box.text_frame
        frame.clear()
        frame.margin_left = frame.margin_right = frame.margin_bottom = 0
        frame.margin_top = 12065
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.NONE
        run = frame.paragraphs[0].add_run()
        run.text = row["en"]
        run.font.name = sample_run.font.name
        run.font.size = sample_run.font.size
        run.font.bold = sample_run.font.bold
        run.font.color.rgb = clone_rgb(sample_run.font.color.rgb)
    presentation.save(OUTPUT_PPTX)


def apply_editorial_overrides(rows):
    for row_number, (russian, english) in EDITORIAL_OVERRIDES.items():
        row = rows[row_number - 1]
        if russian is not None:
            row["ru"] = russian
        row["en"] = english
        row["pause"] = False
    return rows


def main():
    prefix = corrected_prefix()
    boundary = prefix[-1]["ru"]
    remainder = trim_to_boundary(script_units(), boundary)
    compacted = compact_units(remainder)
    translated = translate_rows(compacted)
    rows = prefix + translated
    if not rows[-1]["ru"].lower().startswith("конец"):
        rows.append({"ru": "Конец", "en": "The End", "pause": False})
    apply_editorial_overrides(rows)
    rows = split_overflowing_rows(rows)
    rows = reflow_dangling_boundaries(rows)
    write_outputs(rows)
    print(f"Preserved {len(prefix)} corrected rows and generated {len(translated)} contextual rows; total {len(rows)}")


if __name__ == "__main__":
    main()
