from lxml import html
from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
import copy
import re
from pathlib import Path

ROOT = Path(__file__).resolve().parent
SOURCE_XLSX = ROOT / "Обыкновенное чудо - Субтитры.xlsx"
SOURCE_PPTX = ROOT / "Обыкновенное чудо - Субтитры.pptx"
PILOT_XLSX = ROOT / "Обыкновенное чудо - Субтитры - первые 3 сцены.xlsx"
OUTPUT_XLSX = ROOT / "Обыкновенное чудо - Субтитры - объединенные.xlsx"
OUTPUT_PPTX = ROOT / "Обыкновенное чудо - Субтитры - объединенные.pptx"

ROLES = {
    "ВОЛШЕБНИК", "ЖЕНА ВОЛШЕБНИКА", "МЕДВЕДЬ", "ПРИНЦЕССА", "КОРОЛЬ",
    "МИНИСТР-АДМИНИСТРАТОР", "ПЕРВЫЙ МИНИСТР", "ЭМИЛИЯ", "АМАНДА",
    "ОРИНТИЯ", "ЭМИЛЬ", "ОХОТНИК", "УЧЕНИК", "ПАЛАЧ", "ХОР", "ТРАКТИРЩИК",
}
ALIASES = {
    "ЖЕНА": "ЖЕНА ВОЛШЕБНИКА", "ВОЛШЕБНИЦА": "ЖЕНА ВОЛШЕБНИКА",
    "УЧЕНИК ОХОТНИКА": "УЧЕНИК", "МИНИСТР": "МИНИСТР-АДМИНИСТРАТОР",
    "АДМИНИСТРАТОР": "МИНИСТР-АДМИНИСТРАТОР",
}


def normalized(text):
    return re.sub(r"[^0-9A-Za-zА-Яа-я]+", "", text.replace("ё", "е").replace("Ё", "Е")).lower()


def clean(text):
    text = re.sub(r"🎶.*$", "", text)
    text = text.replace("Волшебник надевает на голову королю корону.", "")
    text = text.replace("Встает, снимает дорожный плащ и отдает Пажу", "")
    text = text.replace("Замахивается на придворных.", "")
    text = re.sub(r"\(+[^)]*\)+", " ", text)
    names = sorted(ROLES | set(ALIASES) | {"ВДВОЕМ"}, key=len, reverse=True)
    text = re.sub(rf"(?:^|\n)\s*(?:{'|'.join(map(re.escape, names))})\s*[.:]\s*", "\n", text, flags=re.I)
    text = re.sub(r"[()]", " ", text)
    text = re.sub(r"\s+([,.;:!?])", r"\1", text)
    text = re.sub(r"\s+", " ", text).strip(" -\t:)")
    text = re.sub(r"^[.;:]{1,2}\s+", "", text)
    return "" if re.fullmatch(r"[.;:…]+", text) else text


def detect_role(text):
    match = re.match(r"^([А-ЯЁ][А-ЯЁ \-]{1,40}?)[.:]\s*(.*)$", text, re.I)
    if not match:
        return None
    name = re.sub(r"\s*-\s*", "-", re.sub(r"\s+", " ", match.group(1).strip().upper()))
    if name not in ROLES and name not in ALIASES:
        return None
    return ALIASES.get(name, name), match.group(2)


def spoken_entries():
    document = html.parse(str(ROOT / "script.html"))
    css = "\n".join(document.xpath("//style/text()"))
    italic_classes = {
        name
        for name, body in re.findall(r"\.(c\d+)\s*\{([^}]*)\}", css)
        if re.search(r"font-style\s*:\s*italic", body, re.I)
    }
    entries = []
    current_role = None
    parts = []
    started = False

    def flush():
        nonlocal current_role, parts
        text = clean(" ".join(parts))
        if current_role and text:
            entries.append({"role": current_role, "text": text})
        current_role = None
        parts = []

    for element in document.xpath("//p|//h1|//h2|//h3|//h4|//li"):
        if "c16" in element.get("class", "").split():
            flush()
            continue
        for span in element.xpath(".//span"):
            if set(span.get("class", "").split()) & italic_classes:
                span.drop_tree()
        text = re.sub(r"\s+", " ", element.text_content().replace("\xa0", " ")).strip()
        if not text:
            flush()
            continue
        if element.tag in {"h1", "h2", "h3", "h4"} or re.match(r"^СЦЕНА\s+\d+\.\d+", text):
            flush()
            continue
        found = detect_role(text)
        if found:
            flush()
            current_role, remainder = found
            if remainder:
                parts.append(remainder)
            started = True
            continue
        role_names = sorted(ROLES | set(ALIASES), key=len, reverse=True)
        inline = re.search(rf"\b(ВДВОЕМ)\s*[.:]\s*|\b({'|'.join(map(re.escape, role_names))})\s*:\s*", text, re.I)
        if not inline:
            inline = re.search(r"\b(ОХОТНИК)\s+(?=А Вы\b)", text, re.I)
        if inline and current_role:
            parts.append(text[:inline.start()])
            flush()
            matched_name = inline.group(1) or (inline.group(2) if inline.lastindex and inline.lastindex > 1 else "")
            name = re.sub(r"\s+", " ", matched_name.upper())
            current_role = "MIXED" if name == "ВДВОЕМ" else ALIASES.get(name, name)
            parts.append(text[inline.end():])
            started = True
            continue
        if not started:
            continue
        if text.startswith("(") and text.endswith(")"):
            flush()
        elif current_role:
            parts.append(text)
    flush()
    return entries


def load_source_rows():
    sheet = load_workbook(SOURCE_XLSX, read_only=True).active
    rows = []
    for index, (_, ru, en) in enumerate(sheet.iter_rows(min_col=1, max_col=3, values_only=True), 1):
        rows.append({
            "source_row": index,
            "ru": str(ru).strip(),
            "en": str(en).strip(),
            "pause": ru == "***",
        })
    return rows


def align_roles(rows, entries):
    content_rows = [row for row in rows[155:-1] if not row["pause"]]
    entry_index = 47
    entry_offset = 0
    for row in content_rows:
        if row["ru"].startswith("Антракт") or row["ru"] in {
            "Первый звонок", "Второй звонок", "Третий звонок", "Министр- администратор:"
        }:
            row["role"] = "MIXED"
            continue
        target = normalized(row["ru"])
        consumed = ""
        assigned = []
        while len(consumed) < len(target) and entry_index < len(entries):
            source = normalized(entries[entry_index]["text"])
            available = source[entry_offset:]
            take = min(len(target) - len(consumed), len(available))
            consumed += available[:take]
            assigned.append(entry_index)
            entry_offset += take
            if entry_offset == len(source):
                entry_index += 1
                entry_offset = 0
        if consumed != target:
            raise ValueError(f"Alignment failed at source row {row['source_row']}: {row['ru']}")
        roles = {entries[index]["role"] for index in assigned}
        row["role"] = next(iter(roles)) if len(roles) == 1 else "MIXED"
    return rows


def fits(rows, dialogue=False):
    english = "\n".join(row["en"] for row in rows) if dialogue else " ".join(row["en"] for row in rows)
    russian = "\n".join(row["ru"] for row in rows) if dialogue else " ".join(row["ru"] for row in rows)
    return len(english) <= (54 if dialogue else 64) and len(english.split()) <= 14 and len(russian.split()) <= 18


def compact(rows):
    result = []
    index = 155
    while index < len(rows):
        row = rows[index]
        if row["pause"] or index == len(rows) - 1:
            result.append({"ru": row["ru"], "en": row["en"], "pause": row["pause"]})
            index += 1
            continue

        group = [row]
        while index + len(group) < len(rows):
            candidate = rows[index + len(group)]
            if row.get("role") in {None, "MIXED"} or candidate["pause"] or candidate.get("role") != row.get("role") or len(group) >= 4:
                break
            if not fits(group + [candidate]):
                break
            group.append(candidate)
        if len(group) > 1:
            result.append({"ru": " ".join(item["ru"] for item in group), "en": " ".join(item["en"] for item in group), "pause": False})
            index += len(group)
            continue

        if index + 1 < len(rows):
            following = rows[index + 1]
            is_exchange = (
                not following["pause"]
                and row.get("role") not in {None, "MIXED"}
                and following.get("role") not in {None, "MIXED", row.get("role")}
                and fits([row, following], dialogue=True)
                and (row["ru"].endswith(("?", "!")) or len(row["ru"].split()) <= 4)
            )
            if is_exchange:
                result.append({
                    "ru": f"— {row['ru']}\n— {following['ru']}",
                    "en": f"— {row['en']}\n— {following['en']}",
                    "pause": False,
                })
                index += 2
                continue
        result.append({"ru": row["ru"], "en": row["en"], "pause": False})
        index += 1
    return result


def pilot_rows():
    sheet = load_workbook(PILOT_XLSX, read_only=True).active
    return [
        {"ru": str(ru).strip(), "en": str(en).strip(), "pause": ru == "***"}
        for _, ru, en in sheet.iter_rows(min_col=1, max_col=3, values_only=True)
    ]


def remove_all_slides(presentation):
    ids = presentation.slides._sldIdLst
    for slide_id in list(ids):
        presentation.part.drop_rel(slide_id.rId)
        ids.remove(slide_id)


def clone_rgb(color):
    return RGBColor(color[0], color[1], color[2])


def write_outputs(rows):
    workbook = load_workbook(SOURCE_XLSX)
    sheet = workbook.active
    ru_style = copy.copy(sheet["B1"]._style)
    en_style = copy.copy(sheet["C1"]._style)
    sheet.delete_rows(1, sheet.max_row)
    for index, row in enumerate(rows, 1):
        sheet.cell(index, 2, row["ru"])._style = copy.copy(ru_style)
        sheet.cell(index, 3, row["en"])._style = copy.copy(en_style)
        if "\n" in row["ru"] or "\n" in row["en"]:
            sheet.row_dimensions[index].height = 31
    workbook.save(OUTPUT_XLSX)

    presentation = Presentation(SOURCE_PPTX)
    sample = next(shape for shape in presentation.slides[0].shapes if shape.has_text_frame)
    background_rgb = presentation.slides[0].background.fill.fore_color.rgb
    sample_run = next(run for paragraph in sample.text_frame.paragraphs for run in paragraph.runs)
    remove_all_slides(presentation)
    layout = min(presentation.slide_layouts, key=lambda item: len(item.placeholders))
    for row in rows:
        slide = presentation.slides.add_slide(layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = clone_rgb(background_rgb)
        if row["pause"]:
            continue
        box = slide.shapes.add_textbox(sample.left, sample.top, presentation.slide_width - 2 * sample.left, 1518285)
        frame = box.text_frame
        frame.clear()
        frame.margin_left = frame.margin_right = frame.margin_bottom = 0
        frame.margin_top = 12065
        frame.word_wrap = True
        frame.auto_size = sample.text_frame.auto_size
        run = frame.paragraphs[0].add_run()
        run.text = row["en"]
        run.font.name = sample_run.font.name
        run.font.size = sample_run.font.size
        run.font.bold = sample_run.font.bold
        run.font.color.rgb = clone_rgb(sample_run.font.color.rgb)
    presentation.save(OUTPUT_PPTX)


def main():
    source = load_source_rows()
    aligned = align_roles(source, spoken_entries())
    rows = pilot_rows() + compact(aligned)
    write_outputs(rows)
    print(f"Compacted {len(source)} source rows to {len(rows)} rows")


if __name__ == "__main__":
    main()
"""


def css_classes_with(css, declaration):
    return {
        name
        for name, body in re.findall(r"\.(c\d+)\s*\{([^}]*)\}", css)
        if re.search(declaration, body, re.IGNORECASE)
    }


def has_class(element, classes):
    return bool(set(element.get("class", [])) & classes)


def fully_italic(element, italic_classes):
    spans = [span for span in element.find_all("span") if span.get_text(" ", strip=True)]
    return bool(spans) and all(has_class(span, italic_classes) for span in spans)


def inline_text(node, italic_classes):
    parts = []
    for child in node.children:
        if isinstance(child, NavigableString):
            parts.append(str(child))
        elif isinstance(child, Tag):
            text = inline_text(child, italic_classes)
            if child.name == "span" and has_class(child, italic_classes):
                core = re.sub(r"^\(+|\)+$", "", text.replace("\xa0", " ").strip()).strip()
                text = f"({core})" if core else ""
            parts.append(text)
    return "".join(parts)


def html_lines():
    soup = BeautifulSoup(SCRIPT_PATH.read_text(encoding="utf-8"), "lxml")
    css = "\n".join(style.get_text() for style in soup.find_all("style"))
    italic_classes = css_classes_with(css, r"font-style\s*:\s*italic")
    song_classes = css_classes_with(css, r"color\s*:\s*#156082") - italic_classes
    lines = []
    for element in soup.select("p, h1, h2, h3, h4, li"):
        cell = element.find_parent("td")
        if cell and not cell.find(["h1", "h2"]):
            continue
        if element.find("a", href=re.compile(r"^#")):
            continue
        italic = fully_italic(element, italic_classes)
        text = element.get_text(" ") if italic else inline_text(element, italic_classes)
        text = re.sub(r"\s+", " ", text.replace("\xa0", " ")).strip()
        if not text:
            lines.append("")
            continue
        if italic and not text.startswith("("):
            text = f"({text})"
        if element.name == "h1":
            text = "\x02" + text
        elif not italic and any(has_class(span, song_classes) for span in element.find_all("span")):
            text = "\x01" + text
        lines.append(text)
    return lines


def role_line(text):
    match = re.match(r"^([А-ЯЁ][А-ЯЁ \-]{1,40}?)[.:]\s*(.*)$", text)
    if not match:
        return None
    role = re.sub(r"\s+", " ", match.group(1).strip().upper())
    if role not in KNOWN_ROLES and role not in ALIASES:
        return None
    return ALIASES.get(role, role), match.group(2)


def clean_text(text):
    text = text.replace("\x01", "")
    text = re.sub(r"\(+[^)]*\)+", " ", text)
    role_names = sorted(KNOWN_ROLES | set(ALIASES) | {"ВДВОЕМ"}, key=len, reverse=True)
    role_pattern = "|".join(re.escape(role) for role in role_names)
    text = re.sub(rf"(?:^|\n)\s*(?:{role_pattern})\s*[.:]\s*", "\n", text, flags=re.IGNORECASE)
    text = re.sub(r"[()]", " ", text)
    text = re.sub(r"\s+([,.;:!?])", r"\1", text)
    text = re.sub(r"\s+", " ", text).strip(" -\t:)")
    text = re.sub(r"^[.;:]{1,2}\s+", "", text)
    return "" if re.fullmatch(r"[.;:…]+", text) else text


def spoken_entries():
    entries = []
    started = False
    role = None
    buffer = []
    song = False

    def flush():
        nonlocal role, buffer, song
        text = clean_text("\n".join(buffer))
        if role and text:
            entries.append({"role": role, "text": text, "song": song})
        role = None
        buffer = []
        song = False

    for raw in html_lines():
        text = raw.strip()
        is_song = text.startswith("\x01")
        text = text.lstrip("\x01")
        if not text:
            flush()
            continue
        if text.startswith("\x02") or re.match(r"^СЦЕНА\s+\d+\.\d+", text):
            flush()
            continue
        detected = role_line(text)
        if detected:
            flush()
            role, remainder = detected
            if remainder:
                buffer.append(("\x01" if is_song else "") + remainder)
            song = is_song
            started = True
            continue
        if not started:
            continue
        if (text.startswith("(") and text.endswith(")")) or text.startswith(("🎵", "🎶")):
            flush()
            continue
        if role:
            buffer.append(("\x01" if is_song else "") + text)
            song = song or is_song
    flush()
    return entries


def normalized(text):
    text = text.replace("ё", "е").replace("Ё", "Е")
    return re.sub(r"[^0-9A-Za-zА-Яа-я]+", "", text).lower()


def source_rows():
    sheet = load_workbook(SOURCE_XLSX, read_only=True).active
    return [
        {"source_row": row, "ru": str(sheet.cell(row, 2).value).strip(), "en": str(sheet.cell(row, 3).value).strip()}
        for row in range(7, sheet.max_row)
        if sheet.cell(row, 2).value != "***"
    ]


def align_rows(rows, entries):
    entry_index = 0
    entry_offset = 0
    for row in rows:
        target = normalized(row["ru"])
        consumed = ""
        assigned = []
        while len(consumed) < len(target) and entry_index < len(entries):
            entry_text = normalized(entries[entry_index]["text"])
            available = entry_text[entry_offset:]
            needed = len(target) - len(consumed)
            take = min(needed, len(available))
            consumed += available[:take]
            assigned.append(entry_index)
            entry_offset += take
            if entry_offset == len(entry_text):
                entry_index += 1
                entry_offset = 0
        if consumed != target:
            raise ValueError(f"Alignment failed at Excel row {row['source_row']}: {row['ru']!r}")
        unique = list(dict.fromkeys(assigned))
        row["role"] = entries[unique[0]]["role"] if len(unique) == 1 else None
        row["song"] = all(entries[index]["song"] for index in unique)
    return rows


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--analyze", action="store_true")
    parser.parse_args()
    entries = spoken_entries()
    rows = align_rows(source_rows(), entries)
    assigned = sum(row["role"] is not None for row in rows)
    print(f"Aligned {len(rows)} Excel rows to {len(entries)} spoken entries; roles assigned to {assigned} rows")
    if assigned != len(rows):
        missing = [row["source_row"] for row in rows if row["role"] is None]
        raise ValueError(f"Rows crossing speaker boundaries: {missing[:20]}")


if __name__ == "__main__":
    main()
"""