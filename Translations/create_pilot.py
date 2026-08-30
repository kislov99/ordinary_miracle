import copy
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
SOURCE_XLSX = ROOT / "Обыкновенное чудо - Субтитры.xlsx"
SOURCE_PPTX = ROOT / "Обыкновенное чудо - Субтитры.pptx"
OUTPUT_XLSX = ROOT / "Обыкновенное чудо - Субтитры - первые 3 сцены.xlsx"
OUTPUT_PPTX = ROOT / "Обыкновенное чудо - Субтитры - первые 3 сцены.pptx"


def plain(*groups):
    return False, groups


def dialogue(*groups):
    return True, groups


SLIDES = [
    plain((1,)), plain((2,)), plain((3,)), plain((4,)), plain((5,)), plain((6,)),

    # Scene 1.1, opening song.
    plain((7, 8)),
    plain((9, 10)),
    plain((11,)),
    plain((12,)),
    plain((13,)),
    plain((14, 15, 16, 17, 18, 19)),
    plain((20,)),
    plain((21,)),

    # Scene 1.2, company song and the Wizard's introduction.
    plain((22,)),
    plain((23, 24)),
    plain((25, 26)),
    plain((27, 28)),
    plain((29, 30)),
    plain((31, 32)),
    plain((33,)),
    plain((34, 35)),
    plain((36,)),

    # Scene 1.3, morning at the Wizard's house.
    plain((37, 38, 39)),
    plain((40, 41, 42, 43)),
    dialogue((44,), (45,)),
    plain((46,)),
    dialogue((47,), (48,)),
    dialogue((49, 50), (51,)),
    plain((52, 53, 54, 55)),
    plain((56,)),
    dialogue((57,), (58,)),
    dialogue((59,), (60,)),
    dialogue((61,), (62,)),
    dialogue((63,), (64,)),
    plain((65, 66)),
    plain((67, 68)),
    plain((69,)),
    plain((70, 71)),
    plain((72, 73)),
    plain((74, 75)),
    plain((76, 77)),
    plain((78,)),
    plain((79, 80)),
    plain((81, 82)),
    dialogue((83, 84), (85,)),
    plain((86, 87)),
    plain((88, 89)),
    plain((90,)),
    plain((91, 92)),
    plain((93,)),
    plain((94, 95)),
    dialogue((96,), (97,)),
    dialogue((98,), (99,)),
    dialogue((100,), (101,)),
    dialogue((102,), (103,)),
    plain((104, 105)),
    plain((106, 107, 108)),
    dialogue((109,), (110,)),
    plain((111, 112)),
    plain((113,)),
    plain((114, 115)),
    plain((116,)),
    plain((117,)),
    dialogue((118,), (119,)),
    dialogue((120,), (121,)),
    plain((122, 123, 124)),
    plain((125, 126)),
    dialogue((127,), (128,)),
    plain((129, 130)),
    plain((131, 132)),
    plain((133, 134)),
    plain((135,)),
    plain((136, 137)),
    plain((138, 139, 140)),
    plain((141, 142)),
    plain((143, 144, 145)),
    plain((146, 147, 148)),
    plain((149, 150)),
    plain((151, 152)),
    dialogue((153,), (154,)),
    plain((155,)),
]

ENGLISH_OVERRIDES = {
    "Приходит день, приходит час, Приходит миг, приходит срок, И рвется связь.":
        "Day, hour, moment, time — then the bond breaks.",
    "Мы воплощение ваших снов, Заветных мыслей, тайных дум!":
        "We embody your dreams, your cherished thoughts and secret hopes!",
    "— Что с тобой? Тебя кто-то обидел?\n— Ты.":
        "— What's wrong? Who hurt you?\n— You.",
    "— Говори, не томи.\n— Ну, что ты натворил нынче утром в курятнике?":
        "— Tell me. Don't stall.\n— What did you do in the coop?",
    "Там вон песок привезли, дорожки посыпать. Превратил бы его в сахар.":
        "They brought sand for the paths. Turn it into sugar.",
    "Ну, тогда те камни, что сложены у амбара, превратил бы в сыр.":
        "Then turn the stones by the barn into cheese.",
    "Иногда нашалишь, а потом все исправишь. А иной раз, щелк! И нет пути назад.":
        "Sometimes a prank can be undone. But sometimes — click! — no way back.",
    "— Я этих цыплят и волшебной палочкой колотил.\n— Угу.":
        "— I tried my magic wand on them.\n— Uh-huh.",
    "И вихрем их завивал! И семь раз молнией их ударил!":
        "I spun them in a whirlwind and struck them with lightning seven times!",
    "Значит, тут уж сделанного не поправишь. Ну что ж, ничего не поделаешь.":
        "So it can't be undone. Well, nothing can be done.",
}


def combine(sheet, column, is_dialogue, groups):
    lines = []
    for group in groups:
        text = " ".join(str(sheet.cell(row, column).value).strip() for row in group)
        lines.append(("— " if is_dialogue else "") + text)
    return "\n".join(lines)


def remove_all_slides(presentation):
    slide_ids = presentation.slides._sldIdLst
    for slide_id in list(slide_ids):
        presentation.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)


def clone_rgb(color):
    return RGBColor(color[0], color[1], color[2])


def main():
    source_workbook = load_workbook(SOURCE_XLSX)
    source_sheet = source_workbook.active
    rows = []
    for is_dialogue, groups in SLIDES:
        russian = combine(source_sheet, 2, is_dialogue, groups)
        rows.append({
            "ru": russian,
            "en": ENGLISH_OVERRIDES.get(russian, combine(source_sheet, 3, is_dialogue, groups)),
        })

    output_workbook = load_workbook(SOURCE_XLSX)
    output_sheet = output_workbook.active
    russian_style = copy.copy(output_sheet["B1"]._style)
    english_style = copy.copy(output_sheet["C1"]._style)
    output_sheet.delete_rows(1, output_sheet.max_row)
    for row_number, row in enumerate(rows, 1):
        output_sheet.cell(row_number, 2, row["ru"])._style = copy.copy(russian_style)
        output_sheet.cell(row_number, 3, row["en"])._style = copy.copy(english_style)
        if "\n" in row["ru"] or "\n" in row["en"]:
            output_sheet.row_dimensions[row_number].height = 31
    output_workbook.save(OUTPUT_XLSX)

    presentation = Presentation(SOURCE_PPTX)
    source_slides = list(presentation.slides)
    selected_slides = []
    for is_dialogue, groups in SLIDES:
        first_source_row = groups[0][0]
        source_slide = source_slides[first_source_row - 1]
        selected_slides.append((source_slide, is_dialogue, groups))

    remove_all_slides(presentation)
    blank_layout = min(presentation.slide_layouts, key=lambda layout: len(layout.placeholders))
    for row, (source_slide, _, _) in zip(rows, selected_slides):
        slide = presentation.slides.add_slide(blank_layout)
        source_background = source_slide.background.fill
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = clone_rgb(source_background.fore_color.rgb)
        if row["ru"] == "***":
            continue
        source_box = next(shape for shape in source_slide.shapes if shape.has_text_frame)
        box = slide.shapes.add_textbox(
            source_box.left,
            source_box.top,
            presentation.slide_width - 2 * source_box.left,
            1518285,
        )
        frame = box.text_frame
        frame.clear()
        frame.margin_left = source_box.text_frame.margin_left
        frame.margin_right = source_box.text_frame.margin_right
        frame.margin_top = source_box.text_frame.margin_top
        frame.margin_bottom = source_box.text_frame.margin_bottom
        frame.word_wrap = True
        frame.auto_size = source_box.text_frame.auto_size
        paragraph = frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = row["en"]
        sample_run = next(
            run
            for paragraph in source_box.text_frame.paragraphs
            for run in paragraph.runs
        )
        run.font.name = sample_run.font.name
        run.font.size = sample_run.font.size
        run.font.bold = sample_run.font.bold
        run.font.color.rgb = clone_rgb(sample_run.font.color.rgb)
    presentation.save(OUTPUT_PPTX)
    print(f"Created {len(rows)} pilot slides from 155 source slides")


if __name__ == "__main__":
    main()